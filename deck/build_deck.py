"""Build the Flipkart Hackathon pitch deck (PPTX) programmatically.
Run: python deck/build_deck.py   -> deck/Karnataka_Traffic_AI.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ───────── theme ─────────
BG          = RGBColor(0x0F, 0x11, 0x15)   # dark navy bg
PANEL       = RGBColor(0x15, 0x18, 0x1F)   # card bg
ACCENT      = RGBColor(0x37, 0x6A, 0xFF)   # primary indigo
ACCENT2     = RGBColor(0x22, 0xC5, 0x5E)   # green
ACCENT3     = RGBColor(0xFF, 0x99, 0x33)   # saffron
TEXT        = RGBColor(0xE6, 0xE7, 0xEA)
TEXT_MUTED  = RGBColor(0x9A, 0xA0, 0xA8)
TEXT_FAINT  = RGBColor(0x6B, 0x72, 0x80)
SAFFRON     = RGBColor(0xFF, 0x99, 0x33)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
INDIA_GREEN = RGBColor(0x13, 0x88, 0x08)
RED         = RGBColor(0xCC, 0x33, 0x33)
BLUE        = RGBColor(0x33, 0x66, 0xCC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def fill_solid(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    fill_solid(bg, color)


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=TEXT, bullet=True, spacing=6, font="Segoe UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    first = True
    for item in items:
        if isinstance(item, tuple):
            label, sub = item
        else:
            label, sub = item, None
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        r = p.add_run()
        r.text = (("• " if bullet else "") + label)
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = font
        r.font.bold = bool(sub)
        if sub:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            p2.space_after = Pt(spacing)
            r2 = p2.add_run()
            r2.text = "    " + sub
            r2.font.size = Pt(size - 2)
            r2.font.color.rgb = TEXT_MUTED
            r2.font.name = font
    return tb


def add_header(slide, title, subtitle=None):
    """Coloured accent bar + title + optional subtitle."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.45), Inches(0.10), Inches(0.55))
    fill_solid(bar, ACCENT)
    add_text(slide, Inches(0.75), Inches(0.40), Inches(12), Inches(0.7), title, size=28, bold=True, color=TEXT)
    if subtitle:
        add_text(slide, Inches(0.75), Inches(1.05), Inches(12), Inches(0.4), subtitle, size=14, color=TEXT_MUTED)


def add_footer(slide, page_num, total):
    add_text(slide, Inches(0.5), Inches(7.0), Inches(6), Inches(0.3),
             "Karnataka Event Traffic AI · Flipkart Hackathon 2026", size=10, color=TEXT_FAINT)
    add_text(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3),
             f"{page_num} / {total}", size=10, color=TEXT_FAINT, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, *, accent=ACCENT, title_size=14, body_size=11):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    fill_solid(card, PANEL)
    card.adjustments[0] = 0.06
    add_text(slide, x + Inches(0.15), y + Inches(0.10), w - Inches(0.3), Inches(0.40),
             title, size=title_size, bold=True, color=accent)
    add_text(slide, x + Inches(0.15), y + Inches(0.55), w - Inches(0.3), h - Inches(0.65),
             body, size=body_size, color=TEXT)
    return card


def add_pill(slide, x, y, text, *, color=ACCENT, fg=WHITE, w=None, size=10):
    if w is None:
        w = Inches(max(0.7, 0.12 * len(text)))
    h = Inches(0.32)
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    fill_solid(p, color)
    p.adjustments[0] = 0.50
    tb = add_text(slide, x, y, w, h, text, size=size, bold=True, color=fg,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return p, w


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT):
    line = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = straight
    line.line.color.rgb = color
    line.line.width = Pt(2)
    # arrowhead
    ln = line.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")


def add_metric_row(slide, x, y, metrics):
    """metrics = [(label, value, color), ...]"""
    n = len(metrics)
    gap_in = 0.15
    total_in = 12.0
    w_in = (total_in - gap_in * (n - 1)) / n
    w = Inches(w_in)
    gap = Inches(gap_in)
    for i, (label, value, c) in enumerate(metrics):
        mx = x + i * (w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mx, y, w, Inches(1.3))
        fill_solid(card, PANEL)
        card.adjustments[0] = 0.07
        add_text(slide, mx, y + Inches(0.15), w, Inches(0.35),
                 label, size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        add_text(slide, mx, y + Inches(0.55), w, Inches(0.7),
                 value, size=28, bold=True, color=c, align=PP_ALIGN.CENTER)


# ───────── slides ─────────
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    slides = []  # list of slide-build functions; counted dynamically
    # We assemble first to know total page count.

    def new_slide():
        s = prs.slides.add_slide(blank)
        add_bg(s)
        return s

    # ── 1. Title ──
    def s1():
        s = new_slide()
        # Tricolour stripe
        stripe_w = Inches(0.18)
        for i, c in enumerate([SAFFRON, WHITE, INDIA_GREEN]):
            r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(i * 2.5), stripe_w, Inches(2.5))
            fill_solid(r, c)
        # Title block
        add_text(s, Inches(1.5), Inches(2.2), Inches(10), Inches(0.5),
                 "AI-Powered Event Traffic Forecasting", size=14, color=ACCENT2)
        add_text(s, Inches(1.5), Inches(2.6), Inches(11), Inches(1.5),
                 "Karnataka Police", size=58, bold=True, color=TEXT)
        add_text(s, Inches(1.5), Inches(4.05), Inches(11), Inches(0.6),
                 "RAG + Agentic LLM + MCP + ISO-27001/GDPR for event-driven congestion mitigation",
                 size=18, color=TEXT_MUTED)
        # Pills
        x0 = Inches(1.5); y0 = Inches(5.0)
        labels = [("CatBoost AUC 0.815", ACCENT),
                  ("8,219-doc FAISS RAG", ACCENT2),
                  ("9-tool MCP server", SAFFRON),
                  ("ISO 27001 + GDPR", RED),
                  ("43 pytest", BLUE)]
        cx = x0
        for t, c in labels:
            _, w = add_pill(s, cx, y0, t, color=c)
            cx += w + Inches(0.2)
        add_text(s, Inches(1.5), Inches(6.7), Inches(10), Inches(0.4),
                 "Flipkart Hackathon 2026 · Bengaluru", size=12, color=TEXT_FAINT)
    s1()

    # ── 2. Problem ──
    def s2():
        s = new_slide()
        add_header(s, "Why event-driven traffic is hard",
                   "Generic time-series + ML models miss the spikes that matter most")
        # left card: problem
        add_card(s, Inches(0.6), Inches(1.7), Inches(6), Inches(5.0),
                 "The problem",
                 "Models trained on historical flow data treat every Tuesday at 6pm the same way.\n\n"
                 "But a cricket match at Chinnaswamy Stadium, a VIP convoy on Mysore Road, a "
                 "monsoon water-logging on Bellary Road, or a sudden protest at Vidhana Soudha "
                 "completely reshape traffic in ways no historical pattern can capture.\n\n"
                 "Traffic officers learn about these events too late — usually when WhatsApp groups "
                 "light up. By then the jam is already an hour deep.",
                 accent=RED, title_size=16, body_size=12)
        # right card: opportunity
        add_card(s, Inches(6.85), Inches(1.7), Inches(6), Inches(5.0),
                 "The opportunity",
                 "Real-time news + social signals + a domain-specific dataset can give the model the "
                 "event awareness it's missing.\n\n"
                 "Combine that with a Retrieval-Augmented Generation pipeline that recalls similar "
                 "past events, an LLM that synthesises an officer-ready briefing, and a planner that "
                 "outputs concrete numbers (officers, barricades, diversion route) — and you get "
                 "actionable mitigation, not just prediction.",
                 accent=ACCENT2, title_size=16, body_size=12)
    s2()

    # ── 3. Solution overview ──
    def s3():
        s = new_slide()
        add_header(s, "Solution overview", "From raw event to officer briefing in one query")
        # pipeline shapes
        boxes = [
            ("Natural-language\nquery", ACCENT),
            ("LLM extracts\nevent attrs", SAFFRON),
            ("CatBoost\npredicts closure", ACCENT2),
            ("Planner →\nofficers/barricades", BLUE),
            ("RAG +\nrealtime context", RED),
            ("LLM writes\nofficer briefing", ACCENT),
        ]
        x = Inches(0.5); y = Inches(2.3); w = Inches(2.0); h = Inches(1.4); gap = Inches(0.13)
        for i, (text, color) in enumerate(boxes):
            bx = x + i * (w + gap)
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y, w, h)
            fill_solid(card, PANEL)
            card.adjustments[0] = 0.10
            add_text(s, bx, y, w, h, text, size=12, bold=True, color=color,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if i < len(boxes) - 1:
                add_arrow(s, bx + w + Emu(0), y + h / 2, bx + w + gap - Emu(0), y + h / 2, color=TEXT_MUTED)
        # below: data feeders
        add_text(s, Inches(0.6), Inches(4.4), Inches(12), Inches(0.4),
                 "Data feeders (continuous)", size=14, bold=True, color=TEXT_MUTED)
        feeders = [
            ("Astram dataset · 8,173 events", BLUE),
            ("Open-Meteo (weather)", ACCENT2),
            ("Google News RSS (India)", SAFFRON),
            ("Reddit RSS (r/bangalore, r/karnataka)", RED),
            ("GDELT 2.0 (geo events)", ACCENT),
            ("OSRM (real diversion routes)", ACCENT2),
            ("data.gov.in (Karnataka)", SAFFRON),
        ]
        x = Inches(0.6); y = Inches(4.85)
        cx = x
        for t, c in feeders:
            _, w = add_pill(s, cx, y, t, color=c, size=10)
            cx += w + Inches(0.15)
            if cx > Inches(12):
                cx = x; y += Inches(0.5)
        # metrics row
        add_metric_row(s, Inches(0.5), Inches(5.85), [
            ("Events trained on", "8,173", ACCENT),
            ("FAISS docs", "8,219", ACCENT2),
            ("Model AUC", "0.815", SAFFRON),
            ("Tests", "43 passing", BLUE),
            ("MCP tools", "9", RED),
        ])
    s3()

    # ── 4. Tech stack ──
    def s4():
        s = new_slide()
        add_header(s, "Tech stack", "All open-source · all production-grade · one Dockerfile")
        groups = [
            ("ML / data",
             ["CatBoost (gradient boosting)", "scikit-learn", "imbalanced-learn (SMOTE-NC)",
              "pandas + numpy", "Optuna (hyperparam search)"]),
            ("RAG / LLM",
             ["sentence-transformers (all-MiniLM-L6-v2)", "FAISS (flat IP index)",
              "Ollama / OpenAI / Anthropic / OpenRouter", "MCP Python SDK"]),
            ("Backend",
             ["FastAPI + Pydantic v2", "Uvicorn", "slowapi (rate limiting)",
              "Prometheus instrumentator", "Fernet · PyJWT · bcrypt · authlib (OIDC)"]),
            ("Frontend",
             ["React 18 + TypeScript", "Vite 5 (proxy + HMR)",
              "Leaflet + react-leaflet (Karnataka map)", "CSS variables for dark/light theme"]),
            ("Data sources",
             ["Astram (anonymised, 8k events)", "Open-Meteo · GDELT · Google News",
              "Reddit RSS · OSRM · data.gov.in (wired)"]),
            ("DevOps",
             ["Multi-stage Dockerfile", "Hugging Face Spaces (Docker SDK, 16 GB)",
              "GitHub Actions · pytest · ruff · pip-audit · npm audit",
              "Git LFS for model + FAISS artefacts"]),
        ]
        cols = 3; rows = 2
        gx = Inches(0.5); gy = Inches(1.8); gw = Inches(4.1); gh = Inches(2.6); gap = Inches(0.15)
        for i, (title, items) in enumerate(groups):
            r = i // cols; c = i % cols
            x = gx + c * (gw + gap); y = gy + r * (gh + gap)
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, gw, gh)
            fill_solid(card, PANEL)
            card.adjustments[0] = 0.07
            add_text(s, x + Inches(0.18), y + Inches(0.12), gw, Inches(0.35),
                     title, size=13, bold=True, color=ACCENT)
            add_bullets(s, x + Inches(0.18), y + Inches(0.50), gw - Inches(0.3), gh - Inches(0.6),
                        items, size=10, color=TEXT, spacing=3)
    s4()

    # ── 5. Astram dataset ──
    def s5():
        s = new_slide()
        add_header(s, "The Astram dataset",
                   "Bengaluru Metropolitan Area · Nov 2023 – Apr 2024 · 150 days · 51 columns")
        # left: dataset stats
        add_metric_row(s, Inches(0.5), Inches(1.8), [
            ("Total events", "8,173", ACCENT),
            ("Causes (taxonomy)", "16", ACCENT2),
            ("Corridors", "30", SAFFRON),
            ("Time span", "150 days", BLUE),
            ("Lat/lon coverage", "100%", RED),
        ])
        # cause table
        add_text(s, Inches(0.5), Inches(3.4), Inches(8), Inches(0.4),
                 "Closure rate varies wildly by cause — strong learning signal",
                 size=14, bold=True, color=TEXT)
        rows = [
            ("vip_movement", "80%", RED),
            ("public_event", "46%", ACCENT3),
            ("protest", "40%", ACCENT3),
            ("tree_fall", "39%", ACCENT3),
            ("construction", "27%", SAFFRON),
            ("procession", "26%", SAFFRON),
            ("water_logging", "8%", ACCENT2),
            ("vehicle_breakdown", "4%", ACCENT2),
            ("accident", "3%", ACCENT2),
            ("pot_holes", "2%", ACCENT2),
        ]
        x = Inches(0.5); y = Inches(3.9); w = Inches(6.5); rh = Inches(0.27)
        for i, (cause, rate, color) in enumerate(rows):
            add_text(s, x, y + i * rh, Inches(2.5), rh, cause, size=11, color=TEXT)
            add_text(s, x + Inches(2.6), y + i * rh, Inches(1.0), rh, rate, size=11, bold=True, color=color)
            # bar
            pct = int(rate.rstrip("%"))
            bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(3.8),
                                     y + i * rh + Inches(0.07),
                                     Inches(2.5 * pct / 80), Inches(0.13))
            fill_solid(bar, color)
        # right card: data quality notes
        add_card(s, Inches(7.5), Inches(3.4), Inches(5.4), Inches(3.3),
                 "Data quality findings",
                 "• Reporting-time bias: events cluster 4-8am and 7pm-midnight, near-zero 9-5pm "
                 "— a logging artifact, not real incident timing. We added an "
                 "is_reporting_window feature to flag this.\n"
                 "• 57% of zone, 70% of junction values are NaN — corridor is the cleanest "
                 "spatial feature.\n"
                 "• end_datetime present for only 4.7% of events — predicting duration directly "
                 "would be too sparse; we predict requires_road_closure instead.",
                 accent=ACCENT2, title_size=14, body_size=11)
    s5()

    # ── 6. CatBoost model ──
    def s6():
        s = new_slide()
        add_header(s, "Forecast model — CatBoost",
                   "Binary classifier on requires_road_closure · class-weighted · 13 features")
        add_metric_row(s, Inches(0.5), Inches(1.8), [
            ("ROC-AUC", "0.815", ACCENT2),
            ("Closure recall", "60%", ACCENT),
            ("Accuracy", "85%", SAFFRON),
            ("Positive base rate", "8.3%", BLUE),
            ("Inference (cpu)", "~3 ms", RED),
        ])
        # left: feature importance
        add_text(s, Inches(0.5), Inches(3.4), Inches(6), Inches(0.4),
                 "Feature importance (top 10)", size=14, bold=True, color=TEXT)
        feats = [
            ("longitude", 13.5), ("police_station", 12.4), ("event_cause", 11.6),
            ("latitude", 10.6), ("hour", 10.4), ("corridor", 9.3),
            ("zone", 8.5), ("dow", 5.9), ("month", 5.8), ("junction", 5.2),
        ]
        x = Inches(0.5); y = Inches(3.9); rh = Inches(0.27)
        for i, (name, pct) in enumerate(feats):
            add_text(s, x, y + i * rh, Inches(2.2), rh, name, size=11, color=TEXT)
            add_text(s, x + Inches(2.2), y + i * rh, Inches(0.7), rh, f"{pct}%",
                     size=11, bold=True, color=ACCENT2)
            bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(2.9),
                                     y + i * rh + Inches(0.07),
                                     Inches(3.0 * pct / 14), Inches(0.13))
            fill_solid(bar, ACCENT2)
        # right: notes
        add_card(s, Inches(7.5), Inches(3.4), Inches(5.4), Inches(3.3),
                 "Why CatBoost over XGBoost",
                 "8 categorical columns dominate this dataset (cause, corridor, junction, "
                 "police_station, zone, priority). CatBoost handles them natively — no "
                 "OrdinalEncoder, no one-hot blowup.\n\n"
                 "Tested SMOTE-NC synthetic oversampling vs class_weights — class_weights "
                 "won on AUC (0.815 vs 0.783) and recall (60% vs 30%). SMOTE-NC trick "
                 "buys precision at recall's cost, which is wrong for ops where missing a "
                 "closure is far worse than a false alarm.",
                 accent=ACCENT, title_size=14, body_size=11)
    s6()

    # ── 7. Planner ──
    def s7():
        s = new_slide()
        add_header(s, "Rule-based planner",
                   "Closure probability → concrete operational outputs")
        # left: inputs
        add_card(s, Inches(0.5), Inches(1.8), Inches(4), Inches(4.8),
                 "Inputs",
                 "• event_cause (e.g. vip_movement)\n"
                 "• corridor (e.g. Mysore Road)\n"
                 "• hour-of-day, weekend flag\n"
                 "• closure_prob (from CatBoost)\n"
                 "• lat / lon",
                 accent=BLUE, body_size=12)
        # middle: rules
        add_card(s, Inches(4.6), Inches(1.8), Inches(4), Inches(4.8),
                 "Rules (tunable)",
                 "expected_crowd = base[cause] × hour_factor × weekend_factor\n\n"
                 "barricades = ⌈prob × cause_severity × 5⌉\n\n"
                 "officers = ⌈crowd / 200⌉ + ⌈barricades / 3⌉\n\n"
                 "severity_score = clip(prob × cause_severity / 3, 0, 1)\n\n"
                 "diversion = nearest corridor centroid (≠ event corridor)",
                 accent=SAFFRON, body_size=11)
        # right: outputs
        add_card(s, Inches(8.7), Inches(1.8), Inches(4.2), Inches(4.8),
                 "Outputs (per event)",
                 "• expected_crowd\n• barricades_needed\n• officers_needed\n"
                 "• severity_score (0–1)\n• diversion_corridor\n"
                 "• diversion_lat, diversion_lon\n\n"
                 "Example — VIP movement on Mysore Road at 18:00 →\n"
                 "  crowd 1,500 · barricades 15 · officers 10 · divert to ORR West 1",
                 accent=ACCENT2, body_size=11)
    s7()

    # ── 8. OSRM routing ──
    def s8():
        s = new_slide()
        add_header(s, "Diversion routes — real driving paths via OSRM",
                   "Not straight lines · 24h-cached · self-hostable")
        # diagram: source event → curved route → destination
        add_text(s, Inches(0.5), Inches(1.9), Inches(6), Inches(0.4),
                 "Before — straight Euclidean line:", size=13, bold=True, color=RED)
        # straight line illustration
        ex = Inches(0.5); ey = Inches(2.4)
        # event pin
        evt = s.shapes.add_shape(MSO_SHAPE.OVAL, ex, ey, Inches(0.35), Inches(0.35))
        fill_solid(evt, RED)
        add_text(s, ex, ey + Inches(0.4), Inches(1.8), Inches(0.3), "Event", size=10, color=TEXT)
        # target
        tx = ex + Inches(5.0); ty = ey + Inches(0.5)
        tgt = s.shapes.add_shape(MSO_SHAPE.OVAL, tx, ty, Inches(0.35), Inches(0.35))
        fill_solid(tgt, ACCENT2)
        add_text(s, tx - Inches(0.4), ty + Inches(0.4), Inches(2.0), Inches(0.3), "Diversion target", size=10, color=TEXT)
        # straight line
        line = s.shapes.add_connector(1, ex + Inches(0.35), ey + Inches(0.18), tx, ty + Inches(0.18))
        line.line.color.rgb = ACCENT2
        line.line.width = Pt(2)
        # after
        add_text(s, Inches(0.5), Inches(4.0), Inches(6), Inches(0.4),
                 "After — real road route from OSRM:", size=13, bold=True, color=ACCENT2)
        evt2 = s.shapes.add_shape(MSO_SHAPE.OVAL, ex, Inches(4.5), Inches(0.35), Inches(0.35))
        fill_solid(evt2, RED)
        tgt2 = s.shapes.add_shape(MSO_SHAPE.OVAL, tx, Inches(5.0), Inches(0.35), Inches(0.35))
        fill_solid(tgt2, ACCENT2)
        # zigzag route
        path_xy = [(Inches(0.85), Inches(4.7)),
                   (Inches(1.6), Inches(4.85)),
                   (Inches(2.3), Inches(4.5)),
                   (Inches(3.0), Inches(4.9)),
                   (Inches(3.8), Inches(4.65)),
                   (Inches(4.5), Inches(5.0)),
                   (Inches(5.3), Inches(4.8))]
        prev = (Inches(0.85), Inches(4.7))
        for pt in path_xy[1:]:
            ln = s.shapes.add_connector(1, prev[0], prev[1], pt[0], pt[1])
            ln.line.color.rgb = ACCENT2
            ln.line.width = Pt(3)
            prev = pt
        # right card
        add_card(s, Inches(7.5), Inches(1.9), Inches(5.4), Inches(4.8),
                 "Why this matters",
                 "Straight lines look right on a map but route officers through buildings, "
                 "lakes, and one-ways.\n\n"
                 "OSRM (Project-OSRM) returns the actual driving polyline — typically "
                 "200-300 GPS points for a Bengaluru corridor.\n\n"
                 "Cached in our SQLite layer for 24 h since corridor centroids barely "
                 "shift over a day.\n\n"
                 "Public demo server by default; production self-hosts via "
                 "OSRM_HOST=http://localhost:5000 (one docker run).\n\n"
                 "If the routing API is unreachable, we fall back to a straight line "
                 "with a dashed style — the map legend shows both.",
                 accent=ACCENT, body_size=11)
    s8()

    # ── 9. RAG ──
    def s9():
        s = new_slide()
        add_header(s, "Retrieval-Augmented Generation",
                   "Ground every LLM briefing in actual precedent")
        # pipeline
        boxes = [
            ("Query\n'protest near MG Road'", ACCENT),
            ("all-MiniLM-L6-v2\nembed → 384 dim", SAFFRON),
            ("FAISS flat IP\n8,219 docs", ACCENT2),
            ("Top-k similar past events\n(score 0.55-0.70)", RED),
        ]
        x = Inches(0.5); y = Inches(1.8); w = Inches(2.95); h = Inches(1.3); gap = Inches(0.13)
        for i, (text, c) in enumerate(boxes):
            bx = x + i * (w + gap)
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y, w, h)
            fill_solid(card, PANEL)
            card.adjustments[0] = 0.10
            add_text(s, bx, y, w, h, text, size=12, bold=True, color=c,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if i < len(boxes) - 1:
                add_arrow(s, bx + w, y + h / 2, bx + w + gap, y + h / 2, color=TEXT_MUTED)
        # corpus composition
        add_text(s, Inches(0.5), Inches(3.5), Inches(12), Inches(0.4),
                 "Corpus mix · 8,219 docs", size=14, bold=True, color=TEXT)
        composition = [
            ("Astram events", 8173, ACCENT),
            ("News articles", 30, ACCENT2),
            ("Reddit posts", 16, SAFFRON),
        ]
        total = sum(c for _, c, _ in composition)
        x = Inches(0.5); y = Inches(4.0); bar_h = Inches(0.5); bar_w = Inches(12)
        cx = x
        for label, count, c in composition:
            w = Inches(bar_w.inches * count / total)
            seg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, y, w, bar_h)
            fill_solid(seg, c)
            add_text(s, cx, y + Inches(0.6), w, Inches(0.3),
                     f"{label} ({count:,})", size=10, color=c, align=PP_ALIGN.CENTER)
            cx += w
        # benefit card
        add_card(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.5),
                 "Why this works for traffic",
                 "Past 'tree_fall on Hennur Main Road' events tell the model what really happened "
                 "last time — average duration, whether closure was needed, what diversion was "
                 "actually used. The LLM weaves those facts into the briefing so the officer sees "
                 "evidence, not just a forecast.",
                 accent=ACCENT, body_size=12)
    s9()

    # ── 10. Orchestrator ──
    def s10():
        s = new_slide()
        add_header(s, "Agentic orchestrator",
                   "One query → six-step pipeline → one operator-ready briefing")
        steps = [
            ("1. Extract", "LLM parses NL query → JSON attrs (cause, corridor, hour, weekend, lat/lon)", SAFFRON),
            ("2. Forecast", "CatBoost predicts closure_prob in ~3 ms", ACCENT2),
            ("3. Plan", "Planner converts prob → crowd, barricades, officers, diversion", BLUE),
            ("4. Retrieve", "FAISS pulls top-3 similar past events", ACCENT),
            ("5. Ground", "Pulls live weather + matching news headlines for the corridor", ACCENT3),
            ("6. Brief", "LLM writes a 4-6 sentence briefing using ALL of the above as context", RED),
        ]
        for i, (label, desc, c) in enumerate(steps):
            y = Inches(1.8 + i * 0.85)
            badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y, Inches(0.6), Inches(0.6))
            fill_solid(badge, c)
            add_text(s, Inches(0.5), y, Inches(0.6), Inches(0.6),
                     str(i + 1), size=18, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, Inches(1.3), y + Inches(0.05), Inches(2.5), Inches(0.4),
                     label.split('. ', 1)[1], size=14, bold=True, color=c)
            add_text(s, Inches(4.0), y + Inches(0.10), Inches(8.5), Inches(0.5),
                     desc, size=12, color=TEXT)
    s10()

    # ── 11. MCP ──
    def s11():
        s = new_slide()
        add_header(s, "Model Context Protocol (MCP) server",
                   "Expose every tool to AI agents (Claude Desktop, custom) · dual-auth")
        # left: 9 tools
        add_card(s, Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.8),
                 "9 tools (stdio + streamable HTTP)",
                 "forecast_closure(event_attrs) → prob\n"
                 "plan_resources(event, prob)   → plan\n"
                 "forecast_and_plan(event)      → combined\n"
                 "retrieve_similar_events(q, k) → RAG\n"
                 "get_weather(lat, lon, hours)  → Open-Meteo\n"
                 "get_news(q, limit)            → Google News\n"
                 "get_reddit(q, limit, days)    → Reddit RSS\n"
                 "get_gdelt(q, limit, timespan) → GDELT\n"
                 "erase_event(event_id)         → GDPR Art 17",
                 accent=ACCENT, body_size=12)
        # right: dual auth
        add_card(s, Inches(6.85), Inches(1.8), Inches(6.0), Inches(4.8),
                 "Dual-authentication model",
                 "Human operators → JWT bearer tokens\n"
                 "  aud=api · roles: viewer / operator / admin\n"
                 "  used by FastAPI dashboard endpoints\n\n"
                 "AI agents → MCP-scoped JWT tokens\n"
                 "  aud=mcp · role=mcp_agent · scoped tool list\n"
                 "  used by MCP HTTP middleware on every call\n\n"
                 "Cross-audience use is rejected at the JWT layer\n"
                 "(InvalidAudienceError). Every tool call is audit-logged with\n"
                 "actor identity and aud claim.",
                 accent=ACCENT2, body_size=12)
    s11()

    # ── 12. Real-time signals ──
    def s12():
        s = new_slide()
        add_header(s, "Real-time data signals",
                   "Five free public sources · TTL-cached · grounded to Karnataka")
        cards = [
            ("Open-Meteo",
             "Hourly forecast at any lat/lon\nNo API key · 1h cache\nFeeds: chat briefing context",
             ACCENT2),
            ("Google News RSS",
             "India locale (hl=en-IN, gl=IN)\n45 articles per query, 30min cache\nAuto-grounded with 'Bengaluru' when query lacks anchor",
             SAFFRON),
            ("Reddit RSS",
             "r/bangalore, r/karnataka, r/india, r/IndiaSpeaks\nWord-boundary keyword filter\n30min cache (RSS, not API — JSON 403s now)",
             RED),
            ("GDELT 2.0",
             "Local news event detection (last 24h)\nKnows about protests, accidents in near-real-time\nFree, no key, 30min cache",
             ACCENT),
            ("data.gov.in",
             "Karnataka accident/traffic datasets\nWired, key-pending\nFails gracefully when key missing",
             BLUE),
            ("BTP advisories",
             "btp.karnataka.gov.in scrape\nJS-rendered table — deferred\nGDELT covers BTP coverage via news",
             TEXT_MUTED),
        ]
        gx = Inches(0.5); gy = Inches(1.8); gw = Inches(4.1); gh = Inches(2.6); gap = Inches(0.15)
        for i, (title, body, c) in enumerate(cards):
            r = i // 3; col = i % 3
            x = gx + col * (gw + gap); y = gy + r * (gh + gap)
            add_card(s, x, y, gw, gh, title, body, accent=c, title_size=13, body_size=11)
    s12()

    # ── 13. LLM flexibility ──
    def s13():
        s = new_slide()
        add_header(s, "LLM backends — runtime switchable",
                   "Local-first by default · cloud fallbacks for production · admin-only swap")
        backends = [
            ("Ollama (local)",
             "Auto-detected via http://localhost:11434\nDefault when reachable\nqwen3:1.7b (1.4 GB) shipped; pull any model from the UI",
             ACCENT2, "Free · Private · GDPR-clean"),
            ("OpenAI",
             "gpt-4o-mini / gpt-4o\nOPENAI_API_KEY in .env or via UI\nUSD per token",
             ACCENT, "Best quality · Cross-border"),
            ("Anthropic",
             "claude-haiku-4-5 / sonnet-4-6 / opus-4-7\nANTHROPIC_API_KEY\nUSD per token",
             SAFFRON, "Strong reasoning"),
            ("OpenRouter",
             "Any model from one key · :free tier available\nllama-3.3-70b-instruct:free, gemini-2.0-flash-exp:free\n0 USD",
             BLUE, "Multi-model · Free tier"),
        ]
        gx = Inches(0.5); gy = Inches(1.8); gw = Inches(6.2); gh = Inches(2.3); gap = Inches(0.15)
        for i, (name, body, c, badge_t) in enumerate(backends):
            row = i // 2; col = i % 2
            x = gx + col * (gw + gap); y = gy + row * (gh + gap)
            add_card(s, x, y, gw, gh, name, body, accent=c, title_size=14, body_size=11)
            add_pill(s, x + gw - Inches(2.5), y + Inches(0.13), badge_t, color=c, fg=WHITE, size=9, w=Inches(2.3))
        # bottom
        add_text(s, Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
                 "Test key → Save settings → runtime swap · No restart · Probe rejects bad keys with the actual provider error", size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    s13()

    # ── 14. Security part 1 (ISO) ──
    def s14():
        s = new_slide()
        add_header(s, "Security & compliance — ISO 27001 Annex A",
                   "Implemented as working code, not documentation theatre")
        controls = [
            ("A.5.1 · Policies",         "SECURITY.md + DPIA.md + RUNBOOK.md",                   ACCENT),
            ("A.8.2 · Classification",   "Every field tagged public / internal / sensitive",     ACCENT2),
            ("A.8.3 · Handling",         "classify.project(row, role) auto-drops fields above clearance", ACCENT2),
            ("A.9.2 · Access mgmt",      "JWT bearer + Depends(require_role('admin'|'operator'|'viewer'))", SAFFRON),
            ("A.9.4 · Source/program",   "Secrets via .env (gitignored) + get_secret() helper",  SAFFRON),
            ("A.10.1 · Cryptography",    "Fernet (AES-128-CBC + HMAC-SHA256) for sensitive fields", RED),
            ("A.12.3 · Backups",         "scripts/backup.{ps1,sh} → tgz of data + models + logs",  BLUE),
            ("A.12.4 · Audit logging",   "Structured JSON per UTC day in logs/audit-*.jsonl",     BLUE),
            ("A.13.1 · Network",         "Caddyfile + docker-compose.prod.yml = HTTPS + HSTS auto-LE", ACCENT),
            ("A.14.2.5 · Secure dev",    "Pydantic v2 validation on every API request",          ACCENT2),
            ("A.16.1 · Incident mgmt",   "RUNBOOK.md: credential compromise, PII leak, model drift",      RED),
            ("A.18.1.4 · Privacy",       "PII regex stripper at ingest (phone, Aadhaar, PAN, ...)", SAFFRON),
        ]
        x = Inches(0.5); y = Inches(1.8); w = Inches(12.3); rh = Inches(0.40)
        for i, (ctl, impl, c) in enumerate(controls):
            row_y = y + i * rh
            add_text(s, x, row_y, Inches(3.5), rh, ctl, size=12, bold=True, color=c)
            add_text(s, x + Inches(3.6), row_y, Inches(8.5), rh, impl, size=11, color=TEXT)
    s14()

    # ── 15. Security part 2 (GDPR) ──
    def s15():
        s = new_slide()
        add_header(s, "GDPR controls",
                   "Privacy by design · operational — not just policy")
        gdpr = [
            ("Art 5(1)(c) · Data minimisation",
             "PII stripper applied at scraper boundary, not after. Classification registry caps field exposure per role.",
             ACCENT),
            ("Art 5(1)(e) · Storage limitation",
             "retention.run() purges audit logs > 365d and scrape cache > 30d. Configurable via env.",
             ACCENT2),
            ("Art 17 · Right to erasure",
             "erasure.erase(event_id) atomically removes from CSV + FAISS index + cache. Admin role + audit-logged.",
             RED),
            ("Art 25 · Privacy by design",
             "PII patterns: phones (+91-aware), Aadhaar, vehicle registration, email, PAN. scan() lets us count residual hits without leaking values.",
             SAFFRON),
            ("Art 30 · Records of processing",
             "Audit log captures actor, action, resource, outcome, custom fields. Append-only, host + pid tagged.",
             ACCENT),
            ("Art 32 · Security of processing",
             "Fernet at rest, JWT in transit, classified access, bcrypt-hashed passwords (≥ 8 chars). Rate-limited endpoints.",
             ACCENT2),
            ("Art 35 · DPIA",
             "DPIA.md documents lawful basis, categories of data, risk register, mitigation, data subject rights.",
             BLUE),
        ]
        x = Inches(0.5); y = Inches(1.8); rh = Inches(0.70)
        for i, (title, body, c) in enumerate(gdpr):
            row_y = y + i * rh
            badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, row_y + Inches(0.07),
                                       Inches(3.5), Inches(0.5))
            fill_solid(badge, c)
            badge.adjustments[0] = 0.30
            add_text(s, x, row_y + Inches(0.07), Inches(3.5), Inches(0.5),
                     title, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, x + Inches(3.7), row_y + Inches(0.10), Inches(9), Inches(0.55),
                     body, size=11, color=TEXT)
    s15()

    # ── 16. Auth model ──
    def s16():
        s = new_slide()
        add_header(s, "Authentication & authorisation",
                   "Three roles · two token types · OIDC-ready")
        roles = [
            ("Viewer",
             "Read-only.\n• See events (public + internal fields, no PII)\n• See state metrics\n• Cannot run forecasts or chat",
             BLUE),
            ("Operator",
             "Everything Viewer + can:\n• Run /forecast, /plan, /chat (LLM)\n• See address, priority, plan details\n• Hit RAG retrieval",
             ACCENT),
            ("Admin",
             "Everything Operator + can:\n• Erase events (GDPR Art 17)\n• Change LLM mode/model/keys at runtime\n• Pull Ollama models · See sensitive fields",
             RED),
        ]
        gx = Inches(0.5); gy = Inches(1.8); gw = Inches(4.1); gh = Inches(2.6); gap = Inches(0.15)
        for i, (name, body, c) in enumerate(roles):
            x = gx + i * (gw + gap)
            add_card(s, x, gy, gw, gh, name, body, accent=c, title_size=15, body_size=11)
        # bottom: token types
        add_card(s, Inches(0.5), Inches(4.55), Inches(6.2), Inches(2.4),
                 "Human JWT (aud=api)",
                 "Issued by /auth/login or /auth/signup\nBcrypt-verified passwords (min 8 chars)\nTTL 1h\nRole in claims, enforced by Depends(require_role(...))\nOIDC scaffold (Google/Microsoft) ready via OIDC_PROVIDER env",
                 accent=ACCENT, body_size=11)
        add_card(s, Inches(6.85), Inches(4.55), Inches(6.0), Inches(2.4),
                 "MCP agent token (aud=mcp)",
                 "Issued by admin via scripts/issue_mcp_token.py\nrole=mcp_agent · scoped tool list in claims\nTTL 24h default\nEnforced by ASGI middleware on the MCP HTTP transport\nCross-audience use rejected with InvalidAudienceError",
                 accent=ACCENT2, body_size=11)
    s16()

    # ── 17. Frontend dashboard ──
    def s17():
        s = new_slide()
        add_header(s, "Operator dashboard",
                   "Karnataka-wide map · live signals · in-line LLM briefing")
        layout = [
            ("Sidebar (left)",
             "• Natural-language query box\n• Filters: cause, corridor, severity slider, date range\n"
             "• State KPIs: events, severe %, total crowd, officers",
             ACCENT),
            ("Map (center)",
             "• Karnataka view (zoom 7) with Leaflet + CartoDB tiles\n"
             "• Severity-coloured circle markers\n"
             "• OSRM diversion polylines (top 25 severe events)\n"
             "• Map legend explaining every symbol",
             ACCENT2),
            ("Right pane",
             "• Live signals (news + Reddit + weather, 60s refresh)\n"
             "• Briefing output: extracted attrs, plan KPIs, officer text\n"
             "• Selected-event details with admin erasure button",
             SAFFRON),
        ]
        gx = Inches(0.5); gy = Inches(1.8); gw = Inches(4.1); gh = Inches(3.5); gap = Inches(0.15)
        for i, (title, body, c) in enumerate(layout):
            x = gx + i * (gw + gap)
            add_card(s, x, gy, gw, gh, title, body, accent=c, title_size=14, body_size=11)
        # bottom features
        feats = [
            ("Dark / Light theme toggle"),
            ("Mobile-responsive (1100 / 768 breakpoints)"),
            ("Karnataka Police branded login with custom SVG emblem"),
            ("Sign up + Login tabs · role chooser · bcrypt-hashed"),
            ("LLM settings modal (admin) with Test key probe"),
            ("Demo accounts auto-fill"),
        ]
        x = Inches(0.5); y = Inches(5.5)
        cx = x
        for t in feats:
            _, w = add_pill(s, cx, y, t, color=ACCENT2, size=10)
            cx += w + Inches(0.15)
            if cx > Inches(12.5):
                cx = x; y += Inches(0.5)
    s17()

    # ── 18. Map legend ──
    def s18():
        s = new_slide()
        add_header(s, "Map legend — every symbol explained",
                   "Operator can interpret the map at a glance")
        # legend mock-up
        x = Inches(0.5); y = Inches(1.8); w = Inches(6.5)
        # title
        add_text(s, x, y, w, Inches(0.4), "LEGEND (as shown in app, bottom-left)",
                 size=11, bold=True, color=TEXT_MUTED)
        # red dot
        rows = [
            (RED, "Event required road closure", "requires_road_closure = True"),
            (BLUE, "Event did not close road", "requires_road_closure = False"),
        ]
        for i, (c, label, sub) in enumerate(rows):
            ry = y + Inches(0.5 + i * 0.5)
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x, ry + Inches(0.10), Inches(0.20), Inches(0.20))
            fill_solid(dot, c)
            add_text(s, x + Inches(0.35), ry, Inches(4), Inches(0.4), label, size=12, color=TEXT)
            add_text(s, x + Inches(0.35), ry + Inches(0.25), Inches(5), Inches(0.3),
                     sub, size=10, color=TEXT_MUTED)
        # severity scale
        sy = y + Inches(1.65)
        add_text(s, x, sy, Inches(6), Inches(0.4), "Size = predicted severity (low → high)",
                 size=12, color=TEXT)
        for i, r in enumerate([0.15, 0.25, 0.40]):
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x + Inches(i * 0.6), sy + Inches(0.45),
                                     Inches(r), Inches(r))
            fill_solid(dot, RED)
            dot.fill.fore_color.rgb = RED
        # diversion routes
        dy = y + Inches(2.85)
        add_text(s, x, dy, Inches(6), Inches(0.4),
                 "Diversion routes", size=12, color=TEXT)
        # solid
        ln1 = s.shapes.add_connector(1, x, dy + Inches(0.6), x + Inches(1.2), dy + Inches(0.6))
        ln1.line.color.rgb = ACCENT2; ln1.line.width = Pt(3)
        add_text(s, x + Inches(1.3), dy + Inches(0.45), Inches(4), Inches(0.3),
                 "Solid · OSRM (real driving path)", size=11, color=TEXT)
        # dashed
        ln2 = s.shapes.add_connector(1, x, dy + Inches(0.95), x + Inches(1.2), dy + Inches(0.95))
        ln2.line.color.rgb = ACCENT2; ln2.line.width = Pt(2)
        ln2.line._get_or_add_ln()
        # dash style
        ln = ln2.line._get_or_add_ln()
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dash")
        add_text(s, x + Inches(1.3), dy + Inches(0.8), Inches(5), Inches(0.3),
                 "Dashed · straight-line fallback (routing unavailable)", size=11, color=TEXT)
        # right card
        add_card(s, Inches(7.5), Inches(1.8), Inches(5.4), Inches(4.8),
                 "Why design legends matter",
                 "Operators glance at the map for ~5 seconds before acting. Without a legend, "
                 "a red dot looks like an alert, not a confirmed closure. A big red circle could "
                 "mean either size = severity or that the event is severe regardless of size.\n\n"
                 "We made the encoding explicit so on-duty officers can read the map without "
                 "training.",
                 accent=ACCENT, body_size=11)
    s18()

    # ── 19. Tests & CI ──
    def s19():
        s = new_slide()
        add_header(s, "Quality — tests, lint, CI",
                   "43 tests · ~22 s · no network deps · GitHub Actions on every push")
        rows = [
            ("test_security.py",    "12 tests", "PII regexes, JWT both-audiences, classification, Fernet, audit"),
            ("test_planner.py",     "7 tests",  "Severity monotonicity, crowd by cause, weekend factor, diversion logic"),
            ("test_data.py",        "6 tests",  "preprocess clean/featurize, cache get/put/expiry/purge"),
            ("test_rag.py",         "3 tests",  "Retrieve top-k, score order, correct cause for known query"),
            ("test_orchestrator.py","1 test",   "End-to-end pipeline with monkeypatched run_llm"),
            ("test_api.py",         "9 tests",  "FastAPI TestClient: roles, password bcrypt, role gating, signup, erasure"),
            ("test_smoke.py",       "1 test",   "Imports every module to catch path/syntax breakage"),
            ("conftest.py",         "fixtures", "Temp AUTH_SECRET, isolated cache, isolated audit log"),
        ]
        x = Inches(0.5); y = Inches(1.8); rh = Inches(0.40)
        for i, (file, count, what) in enumerate(rows):
            ry = y + i * rh
            add_text(s, x, ry, Inches(3), rh, file, size=12, bold=True, color=ACCENT, font="Consolas")
            add_text(s, x + Inches(3.2), ry, Inches(1.5), rh, count, size=11, color=ACCENT2)
            add_text(s, x + Inches(4.8), ry, Inches(8), rh, what, size=11, color=TEXT)
        # CI card
        add_card(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8),
                 "GitHub Actions CI",
                 "On every push / PR — ruff lint · pytest · Docker backend build · Docker frontend build · pip-audit · npm audit\n"
                 "Build matrix pinned to Python 3.11 · Node 20 · ubuntu-latest\n"
                 "Audit jobs are advisory at this maturity (continue-on-error) — flip the toggle when v1.0 ships.",
                 accent=ACCENT2, body_size=11)
    s19()

    # ── 20. Deployment ──
    def s20():
        s = new_slide()
        add_header(s, "Deployment — Hugging Face Spaces (Docker)",
                   "Free · 16 GB RAM · always-on · one container, one URL")
        # diagram
        add_text(s, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
                 "Multi-stage Dockerfile", size=14, bold=True, color=TEXT)
        boxes = [
            ("Stage 1\nnode:20-alpine\nVite build", SAFFRON),
            ("Stage 2\npython:3.11-slim\npip install", ACCENT),
            ("Stage 3 (runtime)\nFastAPI + Uvicorn\n+ React static\n+ FAISS index\n+ CatBoost", ACCENT2),
        ]
        x = Inches(0.5); y = Inches(2.3); w = Inches(4.0); h = Inches(1.8); gap = Inches(0.15)
        for i, (text, c) in enumerate(boxes):
            bx = x + i * (w + gap)
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y, w, h)
            fill_solid(card, PANEL)
            card.adjustments[0] = 0.10
            add_text(s, bx, y, w, h, text, size=12, bold=True, color=c,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if i < len(boxes) - 1:
                add_arrow(s, bx + w, y + h / 2, bx + w + gap, y + h / 2, color=TEXT_MUTED)
        # right side: secrets
        add_card(s, Inches(0.5), Inches(4.4), Inches(6.0), Inches(2.6),
                 "Secrets (HF Settings → Variables and secrets)",
                 "• AUTH_SECRET (JWT signing, ≥ 32 bytes)\n"
                 "• CRYPTO_KEY (Fernet key for at-rest encryption)\n"
                 "• LLM_MODE=openrouter (Ollama not available on HF free)\n"
                 "• OPENROUTER_API_KEY (free $1 credit on signup)\n"
                 "• OPENROUTER_MODEL=meta-llama/llama-3.3-70b-instruct:free",
                 accent=BLUE, body_size=11)
        add_card(s, Inches(6.65), Inches(4.4), Inches(6.2), Inches(2.6),
                 "Why HF Spaces over Render / Vercel / Railway",
                 "Render free (512 MB) OOMs on PyTorch + sentence-transformers.\n"
                 "Vercel serverless can't run 10-15s LLM calls.\n"
                 "Railway works but burns $5 credit in ~2 weeks.\n"
                 "HF Spaces: 16 GB CPU, always-on, free forever — designed for ML demos.",
                 accent=ACCENT, body_size=11)
    s20()

    # ── 21. Live link / Q&A ──
    def s21():
        s = new_slide()
        # big centered URL
        # tricolour top stripe
        for i, c in enumerate([SAFFRON, WHITE, INDIA_GREEN]):
            r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(i * 4.444), 0, Inches(4.444), Inches(0.10))
            fill_solid(r, c)
        add_text(s, Inches(0.5), Inches(0.6), Inches(12), Inches(0.6),
                 "Live now", size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.0),
                 "https://debxorion-flipkart-hackathon-project.hf.space",
                 size=24, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER, font="Consolas")
        add_text(s, Inches(0.5), Inches(2.3), Inches(12), Inches(0.4),
                 "Source code on GitHub", size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        add_text(s, Inches(0.5), Inches(2.7), Inches(12), Inches(0.6),
                 "github.com/deb-X-2005/Flipkart_hackathon_project",
                 size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font="Consolas")
        add_text(s, Inches(0.5), Inches(3.8), Inches(12), Inches(0.4),
                 "Try the demo", size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        # cred bar
        creds = [
            ("admin / admin123",   "full access"),
            ("operator / operator123", "forecasts + chat"),
            ("viewer / viewer123", "read-only"),
        ]
        for i, (u, what) in enumerate(creds):
            x = Inches(2.0 + i * 3.3); y = Inches(4.3)
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.0), Inches(1.0))
            fill_solid(card, PANEL)
            card.adjustments[0] = 0.10
            add_text(s, x, y + Inches(0.10), Inches(3.0), Inches(0.4),
                     u, size=13, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER, font="Consolas")
            add_text(s, x, y + Inches(0.55), Inches(3.0), Inches(0.35),
                     what, size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        # closing
        add_text(s, Inches(0.5), Inches(5.8), Inches(12), Inches(0.6),
                 "Thank you · Questions?", size=24, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        add_text(s, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
                 "Built for Bengaluru Traffic Police · Flipkart Hackathon 2026",
                 size=12, color=TEXT_FAINT, align=PP_ALIGN.CENTER)

    s21()

    # Add footers
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        if i == 1 or i == total:
            continue  # title + thank-you have no footer
        add_footer(slide, i, total)

    out = Path(__file__).parent / "Karnataka_Traffic_AI.pptx"
    prs.save(out)
    print(f"wrote {out} ({out.stat().st_size / 1024:.0f} KB, {total} slides)")


if __name__ == "__main__":
    build()
